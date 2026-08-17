"""Builds the downloadable deck: a cover cut from the client's own anchor
screen, the shift their business is being offered, one slide per product
screen with the AI module on it called out, and a close.

Pure presentation assembly — no AI calls, no cost. Reads whatever the
pipeline already produced (analysis/consult/plan JSON, GeneratedImage rows
and their files on disk) and degrades slide by slide when a piece is
missing rather than failing the export.

Two things make it look like the screens it carries rather than like a
template with pictures dropped in:

  - **A fixed BMV cinematic palette** (dark ground, cyan accent) — the same
    register as the result page, so deck and page read as one deliverable.
    (Palette-sampling from the screens was tried and reverted: it inherits
    whatever hue the screens landed on, and a gold interface made an ochre
    deck.)
  - **Every screen slide states the AI module that is drawn on it**, taken
    from the spec the screen was rendered from (`spec_json`). Those strings
    were sent to the image model as text to render, so the caption under a
    screenshot can be checked against the screenshot. Screens generated
    before the spec was persisted simply get no callout.

Layout lessons this file has already paid for, kept: a composite is
CONTAINed and never cover-cropped (it is a designed frame with its own
margins), a detail crop gets the height its own aspect ratio needs at the
column width, and a closing card is sized from its own text. See
tests/test_deck_layout.py.
"""

import json
import os
import re
from datetime import date

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from app.config import settings
from app.models import Request
from app.pipeline import compositing, screen_story
from app.pipeline._shared import employees_with_ids

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.75)

# One display face for headings, one text face for everything else. Both are
# metric-safe on Windows, macOS and the Linux converters, which a downloaded
# deck has to survive without substitution reflowing the layout.
FONT_DISPLAY = "Segoe UI Light"
FONT_TEXT = "Segoe UI"

_ALERT = RGBColor(0xF8, 0x71, 0x71)
_GOOD = RGBColor(0x4A, 0xDE, 0x80)


def _hex_to_rgb(hex_color: str | None, fallback: str = "#4f46e5") -> RGBColor:
    value = (hex_color or fallback).lstrip("#")
    if len(value) != 6:
        value = fallback.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _abs_image_path(file_path: str) -> str | None:
    """`/uploads/images/{id}/{file}.png` -> absolute path on disk, or None if missing."""
    marker = "/uploads/"
    if marker not in file_path:
        return None
    rel = file_path.split(marker, 1)[1]
    abs_path = os.path.join(settings.UPLOADS_DIR, rel)
    return abs_path if os.path.isfile(abs_path) else None


def _presentation_variant(file_path: str, variant: str) -> str | None:
    """The W4 composite beside a screenshot — `<slug>_0.png` -> `<slug>_hero.png`.

    Returns None when compositing was off or the file predates it, so the
    deck falls back to the raw screenshot rather than losing a slide.
    """
    url = compositing.variant_url(file_path, variant, settings.UPLOADS_DIR)
    return _abs_image_path(url) if url else None


# ── primitives ────────────────────────────────────────────────────────────

def _set_alpha(shape, percent: int) -> None:
    """Makes a solid fill translucent.

    python-pptx has no transparency API, so the alpha element is written
    into the fill's colour directly. Wrapped in a try because it reaches
    past the public surface: a scrim that fails to become translucent
    should cost the slide its dimming, never the whole export.
    """
    try:
        srgb = shape._element.spPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int(percent * 1000))})
        srgb.append(alpha)
    except Exception:  # pragma: no cover — cosmetic only
        pass


def _rect(slide, color: RGBColor, left, top, width, height, *, alpha: int | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(left), int(top), int(width), int(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if alpha is not None:
        _set_alpha(shape, alpha)
    return shape


def _hairline(slide, color: RGBColor, left, top, width, *, alpha: int = 100):
    return _rect(slide, color, left, top, width, Pt(0.75), alpha=alpha)


def _add_bg(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _add_text(
    slide,
    text: str,
    left, top, width, height,
    *,
    size: int,
    color: RGBColor,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = FONT_TEXT,
    anchor=None,
    line_spacing: float | None = None,
    spacing: float | None = None,
):
    box = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = box.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    if spacing:
        # Tracking. Uppercase kickers are unreadable set tight, and letting
        # them breathe is most of what separates this from a default theme.
        try:
            run.font._rPr.set("spc", str(int(spacing * 100)))
        except Exception:  # pragma: no cover — cosmetic only
            pass
    return box


def _kicker(slide, text: str, left, top, color: RGBColor, *, width=Inches(8), size: int = 11):
    return _add_text(
        slide, text.upper(), left, top, width, Inches(0.32),
        size=size, color=color, bold=True, spacing=2.6,
    )


def _add_bullets(
    slide,
    items: list[str],
    left, top, width, height,
    *,
    size: int,
    color: RGBColor,
    accent: RGBColor,
    gap: float = 1.35,
    clamp: int = 0,
):
    box = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(size * gap)
        run = p.add_run()
        # Pain points and outcomes are written as full sentences by a stage
        # with no length budget. Three four-line bullets do not fit a column
        # sized for three two-line ones, and the surplus is drawn straight
        # over the summary strip beneath it.
        run.text = f"›  {_clamp(item, clamp) if clamp else item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = FONT_TEXT
    return box


def _add_accent_bar(slide, color: RGBColor, left, top, width, height=Emu(0)) -> None:
    _rect(slide, color, left, top, width, height or Pt(5))


def _place_image_contain(slide, img_path: str, left, top, width, height):
    """Places an image entirely inside the box, centered, aspect preserved.

    The right choice for a W4 composite: that image is already a designed
    frame with its own margins and shadow, so cropping it to fill a box cuts
    off the presentation the compositor just built.
    """
    with Image.open(img_path) as im:
        native_w, native_h = im.size
    scale = min(width / native_w, height / native_h)
    draw_w, draw_h = int(native_w * scale), int(native_h * scale)
    return slide.shapes.add_picture(
        img_path,
        int(left + (width - draw_w) / 2),
        int(top + (height - draw_h) / 2),
        width=draw_w, height=draw_h,
    )


def _place_image_cover(slide, img_path: str, left, top, width, height):
    """Fills exactly (left, top, width, height) with no distortion and no
    letterboxing — crops the longer axis, like CSS `object-fit: cover`.
    Plain `add_picture(width=, height=)` stretches the source instead.
    """
    with Image.open(img_path) as im:
        native_w, native_h = im.size
    native_ratio = native_w / native_h
    target_ratio = width / height

    pic = slide.shapes.add_picture(img_path, int(left), int(top), width=int(width), height=int(height))
    if native_ratio > target_ratio:
        visible_fraction = target_ratio / native_ratio
        crop = (1 - visible_fraction) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    elif native_ratio < target_ratio:
        visible_fraction = native_ratio / target_ratio
        crop = (1 - visible_fraction) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    return pic


def _clamp(text: str, limit: int) -> str:
    """Trims model prose to what its box can hold, at a sentence end when
    there is one and a word boundary otherwise.

    PowerPoint does not shrink text to fit — it draws it, and the overflow
    lands on whatever is underneath or runs off the slide. A consulting
    summary is the one field with no length discipline at all (the stage
    that writes it is answering a question, not filling a box), and it ran
    four lines into a two-line strip at the foot of the slide.
    """
    text = (text or "").strip()
    if len(text) <= limit:
        return text
    window = text[:limit]
    for end in (". ", "! ", "? "):
        cut = window.rfind(end)
        if cut > limit * 0.55:
            return window[: cut + 1]
    cut = window.rfind(" ")
    return (window[:cut] if cut > 0 else window).rstrip(",;:·• ") + "…"


def _fit(text: str, width, size: int, lines: int = 1) -> str:
    """Clamps text to what fits `lines` lines of a box `width` wide at `size`.

    Character counts were being guessed per call site and kept being wrong by
    a few glyphs, which in a 0.66" strip is the difference between one line
    and a second line that is not there. Derived instead: an average glyph in
    these faces is about half an em, so a line holds width / (0.5 * size/72)
    characters. Deliberately pessimistic — over-trimming loses a word,
    under-trimming overlaps the box beside it.
    """
    return _clamp(text, _fit_chars(width, size, lines))


def _fit_chars(width, size: int, lines: int = 1) -> int:
    """How many characters fit `lines` lines of a box `width` wide at `size`."""
    per_char = Inches(0.5 * size / 72) * 1.18
    return max(8, int(width / per_char) * lines)


def _story_of(img) -> dict | None:
    """The caption for one screen. Shared with the result page so the deck
    and the page never describe the same screenshot differently."""
    return screen_story.from_spec_json(
        getattr(img, "spec_json", None), getattr(img, "role_label", "") or "",
    )


# ── the deck ──────────────────────────────────────────────────────────────

def build_presentation(
    req: Request,
    analysis: dict,
    consult_result: dict,
    plan_result: dict,
    images: list,
) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    concept = plan_result.get("concept_name") or req.business_name

    # Fixed BMV cinematic palette — the same register as the result page.
    # The palette USED to be sampled from the rendered screens (owner call,
    # to stop an indigo frame clashing with dark screenshots), but sampling
    # inherits whatever hue the screens landed on: a gold/brown interface
    # produced a muddy ochre deck (owner verdict on the first fund run:
    # ugly). The screens are dark-cinematic by register, so the site's own
    # dark ground frames any of them; the accent stays BMV cyan everywhere.
    BG = _hex_to_rgb("#020617")
    SURFACE = _hex_to_rgb("#0f172a")
    LINE = _hex_to_rgb("#28364e")
    ACCENT = _hex_to_rgb("#22d3ee")
    ACCENT_SOFT = _hex_to_rgb("#a5f3fc")
    TEXT = _hex_to_rgb("#f8fafc")
    MUTED = _hex_to_rgb("#94a3b8")
    primary = ACCENT  # kept as the name the older slides used

    def chrome(slide, *, label: str = ""):
        """The furniture every interior slide shares: ground, a hairline
        under the header, the running concept name and a slide index."""
        _add_bg(slide, BG)
        _rect(slide, ACCENT, 0, 0, SLIDE_W, Pt(3))
        if label:
            _kicker(slide, label, MARGIN, Inches(0.52), ACCENT)
        _add_text(
            slide, concept.upper(), SLIDE_W - MARGIN - Inches(4), Inches(0.52), Inches(4), Inches(0.32),
            size=10, color=MUTED, align=PP_ALIGN.RIGHT, spacing=2.0,
        )
        # No separate page index. It said "01 / 03" while the kicker beside
        # it already said "SCREEN 01", and the only free corner left was
        # underneath the description. The count goes in the kicker instead.

    # ── Cover: their own screen, dimmed, with the concept over it ─────────
    slide = prs.slides.add_slide(blank)
    _add_bg(slide, BG)
    cover_source = None
    if images:
        # The RAW screenshot here, not the hero composite. The composite is a
        # designed frame on a light backdrop; bled to the slide edges and
        # dimmed it shows two pale margins that read as a rendering fault.
        # The raw screen is edge-to-edge interface, which is what a cinematic
        # cover wants.
        cover_source = _abs_image_path(images[0].file_path) or _presentation_variant(images[0].file_path, "hero")
    if cover_source:
        _place_image_cover(slide, cover_source, 0, 0, SLIDE_W, SLIDE_H)
        # A light wash over the image and a deep band under the type. The
        # step between them is meant to be seen — at 76/88 it read as a seam
        # where a scrim had failed rather than as a deliberate lower third.
        _rect(slide, BG, 0, 0, SLIDE_W, Inches(3.35), alpha=38)
        _rect(slide, BG, 0, Inches(3.35), SLIDE_W, SLIDE_H - Inches(3.35), alpha=93)
        _hairline(slide, ACCENT, 0, Inches(3.35), SLIDE_W, alpha=55)
    _rect(slide, ACCENT, 0, 0, SLIDE_W, Pt(3))

    _kicker(slide, req.business_name, MARGIN, Inches(3.55), ACCENT_SOFT, size=12)
    _add_text(
        slide, concept, MARGIN, Inches(4.0), Inches(11.6), Inches(1.5),
        size=54, color=TEXT, bold=False, font=FONT_DISPLAY,
    )
    _hairline(slide, LINE, MARGIN, Inches(5.62), Inches(11.83))
    _add_text(
        slide, (req.industry or "").strip() or "Bespoke software concept",
        MARGIN, Inches(5.85), Inches(7), Inches(0.4), size=13, color=MUTED,
    )
    _add_text(
        slide, f"Prepared for {req.business_name}  ·  {date.today():%B %Y}",
        SLIDE_W - MARGIN - Inches(6), Inches(5.85), Inches(6), Inches(0.4),
        size=11, color=MUTED, align=PP_ALIGN.RIGHT,
    )

    # ── The shift ─────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    chrome(slide, label="The shift")
    # Two lines of room, because a long concept name makes this headline
    # wrap and PowerPoint will not shrink it — at a one-line box the second
    # line was drawn straight through the panels below.
    _add_text(
        slide, f"What we have — and where {concept} takes you", MARGIN, Inches(1.02),
        Inches(11.8), Inches(1.5), size=30, color=TEXT, font=FONT_DISPLAY,
    )

    # Capped at 3 per column with a wide clearance gap before the summary —
    # a verbose dental brief's 4 long pain points overflowed into it at
    # looser caps, and 3 punchy bullets read better anyway.
    col_w = Inches(5.5)
    left_x, right_x = MARGIN, Inches(7.05)
    top_y, bullets_h = Inches(3.05), Inches(3.05)

    _rect(slide, SURFACE, left_x - Inches(0.3), top_y - Inches(0.45), col_w + Inches(0.6), bullets_h + Inches(0.9), alpha=55)
    _rect(slide, SURFACE, right_x - Inches(0.3), top_y - Inches(0.45), col_w + Inches(0.6), bullets_h + Inches(0.9), alpha=55)
    _rect(slide, _ALERT, left_x - Inches(0.3), top_y - Inches(0.45), Pt(2.5), bullets_h + Inches(0.9))
    _rect(slide, _GOOD, right_x - Inches(0.3), top_y - Inches(0.45), Pt(2.5), bullets_h + Inches(0.9))

    _kicker(slide, "Today", left_x, top_y, _ALERT)
    pain_points = (analysis.get("pain_points") or [])[:3]
    _add_bullets(
        slide, pain_points or ["Manual, reactive operations with no AI layer."],
        left_x, top_y + Inches(0.45), col_w, bullets_h,
        size=14, color=TEXT, accent=_ALERT, gap=1.15, clamp=_fit_chars(col_w, 14, 2),
    )

    _kicker(slide, f"With {concept}", right_x, top_y, _GOOD)
    outcome_lines = [analysis.get("growth_opportunity") or ""] if analysis.get("growth_opportunity") else []
    outcome_lines += (consult_result.get("recommended_features") or [])[:2]
    _add_bullets(
        slide, [x for x in outcome_lines if x][:3],
        right_x, top_y + Inches(0.45), col_w, bullets_h,
        size=14, color=TEXT, accent=_GOOD, gap=1.15, clamp=_fit_chars(col_w, 14, 2),
    )

    _add_text(
        slide, _fit(consult_result.get("consulting_summary") or "", Inches(11.8), 11, lines=2), MARGIN, Inches(6.62),
        Inches(11.8), Inches(0.7), size=11, color=MUTED, line_spacing=1.15,
    )

    # The decompose-era artifacts, read straight off the request row — the
    # deck degrades slide by slide when a piece is missing, same as always.
    modules = json.loads(req.modules_json) if getattr(req, "modules_json", None) else []
    business_case = json.loads(req.business_case_json) if getattr(req, "business_case_json", None) else {}
    playbook = json.loads(req.playbook_json) if getattr(req, "playbook_json", None) else {}

    # ── The blueprint: the product, module by module ──────────────────────
    if modules:
        slide = prs.slides.add_slide(blank)
        chrome(slide, label="The blueprint · part by part")
        _add_text(
            slide, "What we'd actually build", MARGIN, Inches(1.02),
            Inches(11.8), Inches(0.8), size=30, color=TEXT, font=FONT_DISPLAY,
        )
        shown = modules[:6]
        cols = 2
        rows = -(-len(shown) // cols)
        grid_top = Inches(2.1)
        grid_h = Inches(4.9)
        gap = Inches(0.3)
        card_w = (Inches(11.83) - gap * (cols - 1)) / cols
        card_h = (grid_h - gap * (rows - 1)) / rows
        for i, m in enumerate(shown):
            x = MARGIN + (i % cols) * (card_w + gap)
            y = grid_top + (i // cols) * (card_h + gap)
            _rect(slide, SURFACE, x, y, card_w, card_h, alpha=70)
            _rect(slide, ACCENT, x, y, Pt(2.5), card_h)
            _add_text(
                slide, f"{i + 1:02d}", x + Inches(0.25), y + Inches(0.18),
                Inches(0.6), Inches(0.3), size=11, color=ACCENT, bold=True,
            )
            _add_text(
                slide, _fit(m.get("name") or "", card_w - Inches(1.2), 15),
                x + Inches(0.75), y + Inches(0.16), card_w - Inches(1.0), Inches(0.35),
                size=15, color=TEXT, bold=True,
            )
            _add_text(
                slide, _fit(m.get("purpose") or "", card_w - Inches(1.0), 11, lines=2),
                x + Inches(0.75), y + Inches(0.55), card_w - Inches(1.0), Inches(0.6),
                size=11, color=MUTED, line_spacing=1.2,
            )
            spec_ai = ((m.get("spec") or {}).get("ai") or {})
            if spec_ai.get("role"):
                _add_text(
                    slide, _fit(f"AI: {spec_ai['role']}", card_w - Inches(1.0), 10, lines=2),
                    x + Inches(0.75), y + card_h - Inches(0.75), card_w - Inches(1.0), Inches(0.6),
                    size=10, color=ACCENT_SOFT, line_spacing=1.2,
                )

    # ── The blueprint: how this makes money ───────────────────────────────
    if business_case:
        slide = prs.slides.add_slide(blank)
        chrome(slide, label="The blueprint · how this makes money")
        _add_text(
            slide, "The business case", MARGIN, Inches(1.02),
            Inches(11.8), Inches(0.8), size=30, color=TEXT, font=FONT_DISPLAY,
        )
        columns = [
            ("Revenue", [f"{s.get('name')}: {s.get('description')}" for s in (business_case.get("revenue_streams") or [])], _GOOD),
            ("Costs removed", [f"{c.get('cost')}: {c.get('how')}" for c in (business_case.get("costs_removed") or [])], _ALERT),
            ("Pricing levers", list(business_case.get("pricing_levers") or []), ACCENT),
        ]
        col_w2 = Inches(3.75)
        for i, (head, items, tint) in enumerate(columns):
            x = MARGIN + i * (col_w2 + Inches(0.29))
            _rect(slide, SURFACE, x, Inches(2.1), col_w2, Inches(3.7), alpha=60)
            _rect(slide, tint, x, Inches(2.1), col_w2, Pt(2.5))
            _kicker(slide, head, x + Inches(0.25), Inches(2.3), tint, width=col_w2 - Inches(0.5))
            _add_bullets(
                slide, [i for i in items if i][:3] or ["—"],
                x + Inches(0.25), Inches(2.8), col_w2 - Inches(0.5), Inches(2.9),
                size=11, color=TEXT, accent=tint, gap=1.1,
                clamp=_fit_chars(col_w2 - Inches(0.5), 11, 3),
            )
        payback = (business_case.get("payback_logic") or "").strip()
        if payback:
            _add_text(
                slide, _fit(payback, Inches(11.8), 11, lines=2), MARGIN, Inches(6.2),
                Inches(11.8), Inches(0.7), size=11, color=MUTED, line_spacing=1.2,
            )

    # ── One slide per product screen ──────────────────────────────────────
    # The slides the deck exists for: the client seeing their own software,
    # framed, with the AI module on it named.
    for i, img in enumerate(images):
        hero = _presentation_variant(img.file_path, "hero") or _abs_image_path(img.file_path)
        if hero is None:
            continue
        story = _story_of(img)

        slide = prs.slides.add_slide(blank)
        chrome(slide, label=f"Screen {i + 1:02d} / {len(images):02d}")

        # Title and subheading left, the description right, on one band — so
        # the screenshot below gets the full width of the slide.
        _add_text(
            slide, img.role_label or "Product screen", MARGIN, Inches(0.84),
            Inches(5.0), Inches(0.52), size=26, color=TEXT, font=FONT_DISPLAY,
        )
        if story and story["subheading"]:
            _add_text(
                slide, story["subheading"], MARGIN, Inches(1.32), Inches(5.0), Inches(0.3),
                size=11, color=MUTED,
            )
        if story and story["description"]:
            _rect(slide, ACCENT, Inches(6.1), Inches(0.9), Pt(2), Inches(0.76))
            _add_text(
                slide, _fit(story["description"], Inches(6.2), 11, lines=3),
                Inches(6.35), Inches(0.87), Inches(6.23), Inches(0.82),
                size=11, color=MUTED, line_spacing=1.25,
            )

        # ONE image per slide. The IN DETAIL column of two composite crops is
        # gone (owner's call, session 35) — the same instruction the result
        # page got.
        #
        # A 1.6:1 composite on a 16:9 slide is HEIGHT-limited, so every
        # millimetre the header and the caption strip take comes straight off
        # the screenshot. The header was tightened until the picture is the
        # largest this deck has ever drawn it: the two-column layout gave it
        # 7.36" of width, the first one-image pass gave it 7.20" — smaller,
        # which defeated the point of removing the crops.
        hero_top = Inches(1.72)
        hero_h = (Inches(4.85) if story and story["ai"] else Inches(5.35))
        _place_image_contain(slide, hero, MARGIN, hero_top, Inches(11.83), hero_h)

        if story and story["ai"]:
            # The AI module, quoted from the spec the screen was drawn from —
            # so this strip can be checked against the picture above it.
            strip_top = Inches(6.62)
            _rect(slide, SURFACE, MARGIN, strip_top, Inches(11.83), Inches(0.66), alpha=45)
            _rect(slide, ACCENT, MARGIN, strip_top, Pt(2.5), Inches(0.66))
            head = story["ai"]["title"] or "AI on this screen"
            _add_text(
                slide, head.upper(), MARGIN + Inches(0.25), strip_top + Inches(0.09),
                Inches(3.5), Inches(0.28), size=9, color=ACCENT_SOFT, bold=True, spacing=2.0,
            )
            _add_text(
                slide, story["ai"]["headline"], MARGIN + Inches(0.25), strip_top + Inches(0.32),
                Inches(4.6), Inches(0.3), size=12, color=TEXT,
            )
            # Three columns across one 0.66" strip, so each one is given a
            # width that does not reach the next. The rationale box used to
            # run 0.4" into the KPI box and both wrapped to a second line
            # that the strip has no room for.
            tail = "   ·   ".join(x for x in (story["ai"]["rationale"], story["ai"]["confidence"]) if x)
            if tail:
                _add_text(
                    slide, _fit(tail, Inches(3.55), 9), MARGIN + Inches(5.1), strip_top + Inches(0.23),
                    Inches(3.55), Inches(0.3), size=9, color=MUTED,
                )
            if story["tracks"]:
                # One line. KPI labels are short individually and unbounded
                # together; four of them wrapped out of the strip.
                _add_text(
                    slide, _fit("  ·  ".join(story["tracks"][:3]), Inches(2.95), 9),
                    MARGIN + Inches(8.85), strip_top + Inches(0.23), Inches(2.95), Inches(0.3),
                    size=9, color=MUTED, align=PP_ALIGN.RIGHT,
                )

    # ── Technical plan: one slide per module, in the owner's voice ────────
    for mi, m in enumerate(modules[:6]):
        spec = m.get("spec") or {}
        tech = m.get("tech") or {}
        agent = tech.get("ai_agent") or {}
        slide = prs.slides.add_slide(blank)
        chrome(slide, label=f"Technical plan · part {mi + 1:02d} / {min(len(modules), 6):02d}")
        _add_text(
            slide, m.get("name") or "Module", MARGIN, Inches(0.95),
            Inches(11.8), Inches(0.6), size=26, color=TEXT, font=FONT_DISPLAY,
        )
        _add_text(
            slide, _fit(m.get("purpose") or "", Inches(11.8), 12, lines=2), MARGIN, Inches(1.5),
            Inches(11.8), Inches(0.55), size=12, color=MUTED, line_spacing=1.2,
        )

        col_top, col_h = Inches(2.25), Inches(2.55)
        left_w, right_w = Inches(5.6), Inches(5.95)
        right_x2 = MARGIN + left_w + Inches(0.28)

        # Left: what it keeps track of + the screens.
        _rect(slide, SURFACE, MARGIN, col_top, left_w, col_h, alpha=60)
        _rect(slide, ACCENT, MARGIN, col_top, Pt(2.5), col_h)
        _kicker(slide, "What it keeps track of", MARGIN + Inches(0.25), col_top + Inches(0.18), ACCENT_SOFT, width=left_w - Inches(0.5), size=10)
        entities = [e.get("entity") for e in (tech.get("data_model") or []) if e.get("entity")]
        data_line = "  ·  ".join(entities[:6]) or "  ·  ".join((spec.get("data") or [])[:6]) or "—"
        _add_text(
            slide, _fit(data_line, left_w - Inches(0.5), 11, lines=3),
            MARGIN + Inches(0.25), col_top + Inches(0.55), left_w - Inches(0.5), Inches(0.95),
            size=11, color=TEXT, line_spacing=1.25,
        )
        _kicker(slide, "The screens you'll use", MARGIN + Inches(0.25), col_top + Inches(1.55), ACCENT_SOFT, width=left_w - Inches(0.5), size=10)
        _add_text(
            slide, _fit("  ·  ".join((spec.get("screens") or [])[:4]) or "—", left_w - Inches(0.5), 11, lines=2),
            MARGIN + Inches(0.25), col_top + Inches(1.92), left_w - Inches(0.5), Inches(0.6),
            size=11, color=TEXT, line_spacing=1.25,
        )

        # Right: the AI, honestly bounded.
        _rect(slide, SURFACE, right_x2, col_top, right_w, col_h, alpha=60)
        _rect(slide, ACCENT, right_x2, col_top, right_w, Pt(2.5))
        _kicker(slide, "Where the AI works", right_x2 + Inches(0.25), col_top + Inches(0.18), ACCENT_SOFT, width=right_w - Inches(0.5), size=10)
        ai_role = agent.get("purpose") or (spec.get("ai") or {}).get("role") or "No AI in this part — deliberately."
        _add_text(
            slide, _fit(ai_role, right_w - Inches(0.5), 11, lines=2),
            right_x2 + Inches(0.25), col_top + Inches(0.55), right_w - Inches(0.5), Inches(0.65),
            size=11, color=TEXT, line_spacing=1.25,
        )
        never = (agent.get("guardrails") or [])[:1]
        handoff = agent.get("escalation") or (spec.get("ai") or {}).get("hands_off") or ""
        if never:
            _add_text(
                slide, _fit(f"Never: {never[0]}", right_w - Inches(0.5), 10, lines=2),
                right_x2 + Inches(0.25), col_top + Inches(1.3), right_w - Inches(0.5), Inches(0.55),
                size=10, color=_ALERT, line_spacing=1.2,
            )
        if handoff:
            _add_text(
                slide, _fit(f"Hands to you: {handoff}", right_w - Inches(0.5), 10, lines=2),
                right_x2 + Inches(0.25), col_top + Inches(1.9), right_w - Inches(0.5), Inches(0.55),
                size=10, color=MUTED, line_spacing=1.2,
            )

        # Bottom band: built in order + finished when.
        band_top, band_h = Inches(5.0), Inches(1.95)
        _rect(slide, SURFACE, MARGIN, band_top, left_w, band_h, alpha=45)
        _kicker(slide, "Built in this order", MARGIN + Inches(0.25), band_top + Inches(0.15), ACCENT_SOFT, width=left_w - Inches(0.5), size=10)
        steps_raw = [s for s in (tech.get("build_sequence") or [])][:4]
        steps_clean = [re.sub(r"^\s*\d+[.)]\s*", "", s) for s in steps_raw]
        _add_bullets(
            slide, steps_clean or ["—"],
            MARGIN + Inches(0.25), band_top + Inches(0.5), left_w - Inches(0.5), band_h - Inches(0.6),
            size=10, color=TEXT, accent=ACCENT, gap=0.85,
            clamp=_fit_chars(left_w - Inches(0.5), 10, 1),
        )
        _rect(slide, SURFACE, right_x2, band_top, right_w, band_h, alpha=45)
        _kicker(slide, "It's finished when", right_x2 + Inches(0.25), band_top + Inches(0.15), _GOOD, width=right_w - Inches(0.5), size=10)
        _add_bullets(
            slide, (tech.get("done_when") or [])[:3] or ["—"],
            right_x2 + Inches(0.25), band_top + Inches(0.5), right_w - Inches(0.5), band_h - Inches(0.6),
            size=10, color=TEXT, accent=_GOOD, gap=0.9,
            clamp=_fit_chars(right_w - Inches(0.5), 10, 2),
        )

    # ── Your AI team ──────────────────────────────────────────────────────
    employees_slide = (consult_result.get("recommended_ai_employees") or [])[:4]
    if employees_slide:
        slide = prs.slides.add_slide(blank)
        chrome(slide, label="Your AI team")
        _add_text(
            slide, "The employees inside it", MARGIN, Inches(1.02),
            Inches(11.8), Inches(0.8), size=30, color=TEXT, font=FONT_DISPLAY,
        )
        n = len(employees_slide)
        card_w = Inches((11.83 - (n - 1) * 0.25) / n)
        gap = Inches(0.25)
        card_top = Inches(2.3)
        # Sized from the longest "why" rather than fixed — the lesson
        # tests/test_deck_layout.py pins: one-line reasons in a fixed tall
        # box read as content that failed to load, and a pathological reason
        # must not push the card off the slide.
        chars_per_line = max(20, int(card_w / Inches(0.083)))
        longest = max((len(e.get("why", "")) for e in employees_slide), default=0)
        lines = max(1, -(-longest // chars_per_line))
        card_h = min(Inches(3.4), Inches(0.95) + Inches(0.26) * lines)
        for i, emp in enumerate(employees_slide):
            x = MARGIN + i * (card_w + gap)
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(card_top), int(card_w), int(card_h))
            card.fill.solid()
            card.fill.fore_color.rgb = SURFACE
            card.line.color.rgb = LINE
            card.shadow.inherit = False
            tf = card.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.margin_left = Pt(16)
            tf.margin_right = Pt(16)
            tf.margin_top = Pt(18)
            tf.margin_bottom = Pt(18)
            p0 = tf.paragraphs[0]
            r0 = p0.add_run()
            r0.text = emp.get("title", "AI Employee")
            r0.font.bold = True
            r0.font.size = Pt(15)
            r0.font.color.rgb = TEXT
            r0.font.name = FONT_TEXT
            p1 = tf.add_paragraph()
            p1.space_before = Pt(8)
            r1 = p1.add_run()
            r1.text = _clamp(emp.get("why", ""), 260)
            r1.font.size = Pt(11.5)
            r1.font.color.rgb = MUTED
            r1.font.name = FONT_TEXT
            _rect(slide, ACCENT, x, card_top, card_w, Pt(2.5))

    # ── Your playbook ─────────────────────────────────────────────────────
    pb_steps = playbook.get("steps") or []
    if pb_steps:
        slide = prs.slides.add_slide(blank)
        chrome(slide, label="Your playbook")
        _add_text(
            slide, "Every step, in order — and who does it", MARGIN, Inches(1.02),
            Inches(11.8), Inches(0.8), size=30, color=TEXT, font=FONT_DISPLAY,
        )
        who_label = {"you": "YOU", "bmv": "BMV", "partner": "PARTNER"}
        phase_heads = [("before", "Before the build"), ("during", "During the build"), ("after", "After launch")]
        col_w3 = Inches(3.75)
        for i, (phase_id, head) in enumerate(phase_heads):
            x = MARGIN + i * (col_w3 + Inches(0.29))
            steps = [s for s in pb_steps if s.get("phase") == phase_id][:4]
            _rect(slide, SURFACE, x, Inches(2.1), col_w3, Inches(4.5), alpha=60)
            _rect(slide, ACCENT, x, Inches(2.1), col_w3, Pt(2.5))
            _kicker(slide, head, x + Inches(0.25), Inches(2.3), ACCENT_SOFT, width=col_w3 - Inches(0.5))
            lines = [
                f"{who_label.get(s.get('who'), '')} — {s.get('title')}" for s in steps
            ]
            _add_bullets(
                slide, lines or ["—"],
                x + Inches(0.25), Inches(2.8), col_w3 - Inches(0.5), Inches(3.6),
                size=11, color=TEXT, accent=ACCENT, gap=1.15,
                clamp=_fit_chars(col_w3 - Inches(0.5), 11, 2),
            )

    # ── Closing slide ─────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    chrome(slide, label="The next step")
    _add_text(
        slide, f"Ready to make {concept} real?", MARGIN, Inches(1.15), Inches(11.8), Inches(1),
        size=34, color=TEXT, font=FONT_DISPLAY,
    )
    # The two honest paths, stated on the deck the client keeps: this
    # document set is complete enough to execute without us.
    _add_text(
        slide,
        "Two ways forward: we execute this plan for you, module by module — or you take the "
        "blueprint, technical plan and playbook to your own team. They are written to be enough.",
        MARGIN, Inches(1.62), Inches(11.8), Inches(0.6), size=13, color=MUTED, line_spacing=1.25,
    )
    _hairline(slide, LINE, MARGIN, Inches(2.3), Inches(11.83))


    _add_text(
        slide, f"Prepared exclusively for {req.business_name}", MARGIN, Inches(6.75),
        Inches(11.8), Inches(0.5), size=12, color=MUTED,
    )

    return prs


def export_path_for(request_id: int) -> str:
    out_dir = os.path.join(settings.UPLOADS_DIR, "exports")
    os.makedirs(out_dir, exist_ok=True)
    return os.path.join(out_dir, f"{request_id}.pptx")
