"""Pins for the deck's PROPORTIONS — the class of defect this artifact keeps
producing and no test has ever caught.

The pptx has a documented history of distortion and overlap bugs (dd181b8,
d6e8959), and the reason they shipped is that "the picture is squashed" and
"the box is four times taller than its text" are things you see rather than
things that raise. These tests do not replace looking at it — session 33
rebuilt the deck, exported all seven slides through Keynote and read them —
but they hold the specific proportions that looking found wrong, so the
next regression is a red test rather than a slide.

The three detail-crop pins were retired in session 35 when the owner removed
the IN DETAIL column from both the deck and the result page — a test for a
feature that no longer exists is not coverage. What replaced them is the new
rule: one image per screen slide, given the width the column freed.

Deliberately NOT pinned here: anything about how the deck looks that a
number cannot express. That still requires
`python scripts/deck_sample.py` and a pair of eyes.
"""

import io
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import pytest
from PIL import Image
from pptx.util import Inches

from app.config import settings
from app.pipeline import export_pptx


class _Row:
    def __init__(self, role_id, role_label, file_path):
        self.role_id, self.role_label, self.variant, self.file_path = role_id, role_label, 0, file_path


class _Req:
    business_name = "Northgate Coffee Roasters"
    business_description = "A coffee roastery."
    industry = "Coffee Roastery"


def _png(path, size):
    Image.new("RGB", size, "white").save(path, format="PNG")


@pytest.fixture
def run_dir(tmp_path, monkeypatch):
    """A screen with its two W4 composites, at the aspect ratios the
    compositor really produces: a wide thin top band and a squarer content
    block."""
    images = tmp_path / "images" / "1"
    images.mkdir(parents=True)
    _png(images / "analytics_0.png", (1408, 814))
    _png(images / "analytics_hero.png", (1548, 957))
    _png(images / "analytics_detail_1.png", (1317, 358))   # 3.68:1
    _png(images / "analytics_detail_2.png", (1007, 522))   # 1.93:1
    monkeypatch.setattr(settings, "UPLOADS_DIR", str(tmp_path))
    return tmp_path


def _build(employees):
    return export_pptx.build_presentation(
        _Req(),
        {"pain_points": ["a", "b"], "growth_opportunity": "c"},
        {"consulting_summary": "s", "recommended_features": ["f"], "recommended_ai_employees": employees},
        {"concept_name": "Northgate Roast Intelligence", "roles": [], "visual_theme": {"primary_color": "#7c3f21"}},
        [_Row("analytics", "Analytics", "/uploads/images/1/analytics_0.png")],
    )


def _pictures(slide):
    return [s for s in slide.shapes if s.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE


def _screen_slide(prs):
    """The product-screen slide, found by what is ON it rather than by index.

    It used to be `prs.slides[1]`. Adding a cover slide silently shifted
    every one of these assertions onto the wrong slide, where they still
    passed by finding nothing.
    """
    for slide in prs.slides:
        text = " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
        if "SCREEN 01" in text and _pictures(slide):
            return slide
    raise AssertionError("no product-screen slide")


def _employee_cards(slide):
    """The rounded cards, not the full-width accent bar — both are autoshapes."""
    return [s for s in slide.shapes if s.name.startswith("Rounded Rectangle")]


def _team_slide(prs):
    """The AI-team slide, found by what is ON it rather than by index — the
    cards moved off the closing slide onto their own slide, and an index
    would silently assert against the wrong one (same lesson as
    _screen_slide)."""
    for slide in prs.slides:
        if _employee_cards(slide):
            return slide
    raise AssertionError("no slide carries employee cards")


def test_a_screen_slide_carries_exactly_one_image(run_dir):
    """The owner's rule (session 35): the screenshot is the subject, and the
    IN DETAIL column of two composite crops is gone from both the deck and
    the result page. Two extra thumbnails competing with the screen is the
    clutter this replaced."""
    pictures = _pictures(_screen_slide(_build([{"title": "AI Ops", "why": "short"}])))

    assert len(pictures) == 1, f"a screen slide should show one image, found {len(pictures)}"


def test_the_screen_takes_the_space_the_crop_column_freed(run_dir):
    """Losing the crop column has to buy the screenshot that space — otherwise
    the slide is just emptier.

    A 16:9-ish screenshot on a 7.5" slide with a header and a caption strip is
    HEIGHT-limited, not width-limited, so "spans the slide" is the wrong
    measure: it fills the vertical band it is given, and it is centred in the
    width rather than parked in the left two thirds.
    """
    picture = _pictures(_screen_slide(_build([{"title": "AI Ops", "why": "short"}])))[0]

    assert picture.height >= Inches(4.4), "the screen no longer fills its band"
    centre = picture.left + picture.width / 2
    assert centre == pytest.approx(export_pptx.SLIDE_W / 2, abs=Inches(0.1)), (
        "the screen is still offset for a two-column slide"
    )
    assert picture.left >= 0 and picture.left + picture.width <= export_pptx.SLIDE_W
    assert picture.top + picture.height <= export_pptx.SLIDE_H


def test_a_closing_card_is_sized_from_its_own_text(run_dir):
    """Fixed at 2.9", two employees with one-line reasons produced two boxes
    four times taller than their content — which reads as content that
    failed to load."""
    short = _build([{"title": "AI Ops", "why": "Short."}])
    long = _build([{"title": "AI Ops", "why": "A much longer sentence about what this employee does all day, " * 3}])

    def _card_height(prs):
        cards = _employee_cards(_team_slide(prs))
        assert cards, "the team slide has no employee cards"
        return cards[0].height

    assert _card_height(long) > _card_height(short)
    assert _card_height(short) < Inches(2.0), "a one-line reason does not need a two-inch box"


def test_the_card_never_grows_without_bound(run_dir):
    """A pathological `why` must not push the card off the slide."""
    prs = _build([{"title": "AI Ops", "why": "x" * 4000}])
    cards = _employee_cards(_team_slide(prs))
    assert cards[0].top + cards[0].height < export_pptx.SLIDE_H
