"""Pins for the deck's two session-35 changes: it takes its colour from the
screens it carries, and it states the AI module drawn on each one.

The palette tests are all about the guards rather than the extraction. A
palette read out of an arbitrary image has no obligation to be usable, and
an illegible slide is worse than a generic one — so what is pinned is that
every path returns a complete, contrast-checked palette, including the
paths where the images tell us nothing.

The callout tests pin the same rule the result page follows: an AI panel is
described only when the spec says one was drawn. Advertising AI on a screen
that does not show it is the one failure here that would cost a client's
trust rather than a slide.
"""

import json
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import pytest
from PIL import Image
from pptx.util import Inches

from app.config import settings
from app.pipeline import deck_palette, export_pptx


def _write(path, color, size=(320, 180)):
    Image.new("RGB", size, color).save(path, format="PNG")
    return str(path)


# ── the palette ───────────────────────────────────────────────────────────

def test_the_accent_is_always_legible_on_the_surface_it_sits_on(tmp_path):
    """A dark navy screen with a deep, barely-visible accent still has to
    produce a readable deck."""
    path = _write(tmp_path / "s.png", (10, 14, 22))
    Image.open(path)
    im = Image.new("RGB", (320, 180), (10, 14, 22))
    im.paste(Image.new("RGB", (120, 60), (18, 40, 70)), (20, 20))
    im.save(path)

    palette = deck_palette.from_images([path], "#0e9594")
    ratio = deck_palette.contrast(
        deck_palette._rgb(palette["accent"]), deck_palette._rgb(palette["surface"]),
    )
    assert ratio >= 4.5, f"accent {palette['accent']} is unreadable on {palette['surface']}"


def test_body_text_always_clears_the_ground(tmp_path):
    path = _write(tmp_path / "s.png", (240, 240, 245))  # a light screen

    palette = deck_palette.from_images([path], None)
    assert deck_palette.contrast(
        deck_palette._rgb(palette["text"]), deck_palette._rgb(palette["bg"]),
    ) >= 7


def test_a_greyscale_screen_borrows_the_brand_rather_than_inventing_an_accent(tmp_path):
    """There is genuinely no accent in a grey image; picking the most
    'colourful' grey would paint the deck in noise."""
    path = _write(tmp_path / "grey.png", (90, 90, 90))

    palette = deck_palette.from_images([path], "#0e9594")
    assert palette["source"].startswith("brand")


def test_an_unreadable_file_never_fails_the_export(tmp_path):
    palette = deck_palette.from_images([str(tmp_path / "does-not-exist.png")], None)

    assert palette["source"] == "fallback"
    assert set(palette) >= {"bg", "surface", "line", "accent", "text", "muted"}


def test_every_path_returns_a_complete_palette(tmp_path):
    keys = {"bg", "surface", "line", "accent", "accent_soft", "text", "muted", "source"}
    for palette in (
        deck_palette.from_images([], None),
        deck_palette.from_images([], "#7c3f21"),
        deck_palette.from_images([_write(tmp_path / "a.png", (12, 30, 40))], "#7c3f21"),
    ):
        assert set(palette) == keys
        assert all(v.startswith("#") for k, v in palette.items() if k != "source")


# ── the screen callout ────────────────────────────────────────────────────

class _Row:
    def __init__(self, role_id, role_label, file_path, spec_json=None):
        self.role_id, self.role_label, self.variant = role_id, role_label, 0
        self.file_path, self.spec_json = file_path, spec_json


class _Req:
    business_name = "Harbourline Marine"
    business_description = "A boatyard."
    industry = "Boatyard"


@pytest.fixture
def run_dir(tmp_path, monkeypatch):
    images = tmp_path / "images" / "1"
    images.mkdir(parents=True)
    _write(images / "dashboard_0.png", (12, 30, 40), (1408, 814))
    monkeypatch.setattr(settings, "UPLOADS_DIR", str(tmp_path))
    return tmp_path


def _build(spec_json):
    return export_pptx.build_presentation(
        _Req(),
        {"pain_points": ["a"], "growth_opportunity": "b"},
        {"consulting_summary": "s", "recommended_features": ["f"],
         "recommended_ai_employees": [{"title": "AI Ops", "why": "short"}]},
        {"concept_name": "Harbourline Navigator", "roles": [], "visual_theme": {"primary_color": "#005f73"}},
        [_Row("dashboard", "Dashboard", "/uploads/images/1/dashboard_0.png", spec_json)],
    )


def _all_text(prs) -> str:
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
    return "\n".join(out)


def test_the_slide_quotes_the_ai_module_the_screen_was_drawn_with(run_dir):
    text = _all_text(_build(json.dumps({
        "subheading": "Today's Operational Snapshot",
        "kpis": [{"label": "Active Jobs"}],
        "ai": {"title": "Optimize Travel Lift 2", "headline": "High demand, low current use",
               "rationale": "Reallocate two slots", "confidence": "92% optimal"},
    })))

    assert "OPTIMIZE TRAVEL LIFT 2" in text
    assert "High demand, low current use" in text
    assert "92% optimal" in text
    assert "Today's Operational Snapshot" in text


def test_no_ai_module_drawn_means_no_ai_claim_on_the_slide(run_dir):
    """An AI panel with no headline was never rendered."""
    text = _all_text(_build(json.dumps({"subheading": "Snapshot", "ai": {"title": "AI Insights", "headline": "  "}})))

    assert "AI Insights" not in text
    assert "AI on this screen" not in text


def test_a_screen_from_before_the_spec_was_kept_says_nothing_about_ai(run_dir):
    text = _all_text(_build(None))

    assert "AI on this screen" not in text


def test_an_unparseable_spec_never_fails_the_export(run_dir):
    assert _build("{not json").slides is not None


# ── overflow, the defect class this file keeps producing ──────────────────

def test_model_prose_is_trimmed_to_what_its_box_can_hold():
    """PowerPoint does not shrink text to fit — it draws it over whatever is
    underneath. The consulting summary ran four lines into a two-line strip."""
    long = ("To overcome manual inefficiencies and unlock significant growth, the business should "
            "implement an AI-powered operational platform. It will centralize scheduling and "
            "eliminate revenue loss across every part of the yard.")
    trimmed = export_pptx._clamp(long, 210)

    assert len(trimmed) <= 211
    assert trimmed.endswith((".", "…")), "a trim should land on a sentence or an ellipsis"


def test_a_short_string_is_left_exactly_as_written():
    assert export_pptx._clamp("Short enough.", 210) == "Short enough."


def test_the_headline_gets_two_lines_before_the_panels_start(run_dir):
    """A long concept name wraps this headline; at a one-line box the second
    line was drawn straight through the panels below it."""
    prs = _build(None)
    shift = prs.slides[1]
    titles = [s for s in shift.shapes if s.has_text_frame and "What we have" in s.text_frame.text]
    assert titles, "the shift slide lost its headline"
    assert titles[0].height >= Inches(1.4)
